#!/usr/bin/env python3
"""slides-premium v2 — 16:9 widescreen, fonts ×1.5, better layout, real LINE data."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

W = Inches(13.33)
H = Inches(7.5)

SAGE     = RGBColor(0x87,0xA0,0x80); SAGE_L = RGBColor(0xEC,0xF2,0xEB)
SAGE_D   = RGBColor(0x55,0x72,0x52); CREAM  = RGBColor(0xF8,0xF5,0xEF)
DARK     = RGBColor(0x2A,0x2A,0x2A); MID    = RGBColor(0x5E,0x5E,0x5E)
WHITE    = RGBColor(0xFF,0xFF,0xFF); RED    = RGBColor(0xDC,0x26,0x26)
GGOOD    = RGBColor(0x16,0x65,0x34); GBORD  = RGBColor(0x22,0xC5,0x5E)
GBG      = RGBColor(0xF0,0xFD,0xF4)
WARN_BG  = RGBColor(0xFF,0xF9,0xE6); WARN_FG= RGBColor(0x7A,0x55,0x00)
WARN_BD  = RGBColor(0xCC,0xA0,0x10)
HINT_BG  = RGBColor(0xEE,0xF2,0xFF); HINT_FG= RGBColor(0x37,0x30,0xA3)
HINT_BD  = RGBColor(0x63,0x66,0xF1)
DB_BG    = RGBColor(0x0F,0x17,0x2A); DB_MID = RGBColor(0x1E,0x29,0x3B)
DB_GRN   = RGBColor(0x3E,0xCF,0x8E); DB_RED = RGBColor(0xF8,0x71,0x71)
DB_TXT   = RGBColor(0xCB,0xD5,0xE1); DB_MUT = RGBColor(0x94,0xA3,0xB8)
CODE_BG  = RGBColor(0x1E,0x29,0x3B); CODE_BL= RGBColor(0x7D,0xD3,0xFC)
CODE_GN  = RGBColor(0x86,0xEF,0xAC); CODE_PR= RGBColor(0xC0,0x84,0xFC)
CODE_GR  = RGBColor(0x47,0x55,0x69); CODE_YL= RGBColor(0xFC,0xD3,0x4D)
LINE_GN  = RGBColor(0xEE,0xFF,0xDD); LINE_HD= RGBColor(0x36,0x3D,0x44)

prs = Presentation()
prs.slide_width = W; prs.slide_height = H
BL = prs.slide_layouts[6]

# ── helpers ─────────────────────────────────────────────────────────────────

def S():
    return prs.slides.add_slide(BL)

def BG(sl, c):
    f = sl.background.fill; f.solid(); f.fore_color.rgb = c

def BOX(sl, x,y,w,h, fill=None, border=None, bw=Pt(1.5)):
    sh = sl.shapes.add_shape(1,x,y,w,h)
    sh.fill.solid() if fill else sh.fill.background()
    if fill: sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = bw
    else: sh.line.fill.background()
    return sh

def T(sl, txt, x,y,w,h, sz=Pt(18), bold=False, c=None, a=PP_ALIGN.LEFT, wrap=True):
    c = c or DARK
    tb = sl.shapes.add_textbox(x,y,w,h); tf = tb.text_frame; tf.word_wrap=wrap
    p = tf.paragraphs[0]; p.alignment = a; r = p.add_run()
    r.text=txt; r.font.size=sz; r.font.bold=bold; r.font.color.rgb=c
    return tb

def MT(sl, lines, x,y,w,h, a=PP_ALIGN.LEFT):
    """lines = [(text, Pt, bold, RGBColor)]"""
    tb = sl.shapes.add_textbox(x,y,w,h); tf = tb.text_frame; tf.word_wrap=True
    first=True
    for ln in lines:
        t,sz,b,cl = ln
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.alignment = a; r = p.add_run()
        r.text=t; r.font.size=sz; r.font.bold=b; r.font.color.rgb=cl
    return tb

def BADGE(sl, txt, x,y, w=Inches(2.2)):
    b = BOX(sl,x,y,w,Inches(0.4),fill=SAGE)
    tf=b.text_frame; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=txt; r.font.size=Pt(13); r.font.bold=True; r.font.color.rgb=WHITE

def HDR(sl, badge, title, bg=WHITE):
    BG(sl,bg); BADGE(sl,badge,Inches(0.5),Inches(0.28))
    T(sl,title,Inches(0.5),Inches(0.75),Inches(12.3),Inches(0.7),sz=Pt(30),bold=True)
    BOX(sl,Inches(0.5),Inches(1.55),Inches(12.3),Inches(0.05),fill=SAGE_L)

def TIP(sl, title, body, x,y,w,h, bg=SAGE_L, tc=SAGE_D, bc=SAGE):
    BOX(sl,x,y,Inches(0.1),h,fill=bc)
    b=BOX(sl,x+Inches(0.1),y,w-Inches(0.1),h,fill=bg)
    b.line.fill.background()
    tf=b.text_frame; tf.word_wrap=True
    p0=tf.paragraphs[0]; r0=p0.add_run()
    r0.text=title; r0.font.size=Pt(16); r0.font.bold=True; r0.font.color.rgb=tc
    p1=tf.add_paragraph(); r1=p1.add_run()
    r1.text=body; r1.font.size=Pt(14); r1.font.color.rgb=DARK

def CARD(sl, icon, title, body, x,y,w,h):
    b=BOX(sl,x,y,w,h,fill=SAGE_L); b.line.fill.background()
    tf=b.text_frame; tf.word_wrap=True
    p0=tf.paragraphs[0]; r0=p0.add_run()
    r0.text=icon+"  "+title; r0.font.size=Pt(17); r0.font.bold=True; r0.font.color.rgb=DARK
    p1=tf.add_paragraph(); r1=p1.add_run()
    r1.text=body; r1.font.size=Pt(13); r1.font.color.rgb=MID

def CIRCLE(sl, n, x,y):
    b=BOX(sl,x,y,Inches(0.5),Inches(0.5),fill=SAGE); b.line.fill.background()
    tf=b.text_frame; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=str(n); r.font.size=Pt(18); r.font.bold=True; r.font.color.rgb=WHITE

def STEP(sl, n, title, body, y, xo=Inches(0.5)):
    CIRCLE(sl,n,xo,y+Inches(0.05))
    T(sl,title,xo+Inches(0.65),y,Inches(11.8),Inches(0.45),sz=Pt(18),bold=True)
    if body:
        T(sl,body,xo+Inches(0.65),y+Inches(0.48),Inches(11.8),Inches(0.42),sz=Pt(14),c=MID)

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — COVER
# ════════════════════════════════════════════════════════════════════════════
sl = S(); BG(sl,SAGE)
T(sl,"สไลด์สอนงาน · แพ็ก Premium",Inches(2),Inches(1.0),Inches(9.33),Inches(0.55),
  sz=Pt(14),bold=True,c=WHITE,a=PP_ALIGN.CENTER)
T(sl,"🚀",Inches(5.66),Inches(1.65),Inches(2),Inches(0.9),
  sz=Pt(52),c=WHITE,a=PP_ALIGN.CENTER)
T(sl,"สอนใช้งาน แพ็ก Premium",Inches(1),Inches(2.65),Inches(11.33),Inches(1.0),
  sz=Pt(54),bold=True,c=WHITE,a=PP_ALIGN.CENTER)
T(sl,"QR Code เมนูครบวงจร  ·  ออเดอร์เข้า LINE OA  ·  Delivery  ·  2 ภาษา",
  Inches(1.5),Inches(3.8),Inches(10.33),Inches(0.55),sz=Pt(18),c=WHITE,a=PP_ALIGN.CENTER)
hb=BOX(sl,Inches(4.16),Inches(4.6),Inches(5),Inches(0.6),fill=RGBColor(0xFF,0xFF,0xFF))
hb.line.color.rgb=WHITE; hb.line.width=Pt(1)
tf2=hb.text_frame; p2=tf2.paragraphs[0]; p2.alignment=PP_ALIGN.CENTER; r2=p2.add_run()
r2.text="← กด ← → หรือปัดซ้าย/ขวา →"; r2.font.size=Pt(14); r2.font.color.rgb=SAGE_D

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — ฟีเจอร์
# ════════════════════════════════════════════════════════════════════════════
sl = S(); HDR(sl,"ภาพรวม","แพ็ก Premium — ได้อะไรบ้าง?")
feats=[
    ("💬","ออเดอร์เข้า LINE OA","ลูกค้ากดยืนยัน → ร้านได้รับออเดอร์ครบใน LINE ทันที ไม่มีจดผิด"),
    ("🏠","Delivery ส่งถึงบ้าน","กรอกที่อยู่ + ชำระ PromptPay ครบใน 1 ขั้นตอน"),
    ("🌐","2 ภาษา ไทย / EN","กดปุ่มสลับ ทุกหน้าเปลี่ยนเป็นอังกฤษพร้อมกัน"),
    ("🎛️","ตัวเลือกเมนูครบ","ขนาด · ประเภทเนื้อ · ท็อปปิ้ง · ความหวาน ราคาปรับอัตโนมัติ"),
    ("✏️","แก้ราคาเองได้","Supabase Dashboard — คลิก แก้ Enter เสร็จ ไม่ต้องรอเรา"),
    ("👁️","ซ่อน/แสดงเมนู","ของหมด — ซ่อนได้ 1 คลิก โดยไม่ต้องลบออกถาวร"),
]
cols=[(Inches(0.5),Inches(1.7)),(Inches(4.61),Inches(1.7)),(Inches(8.72),Inches(1.7)),
      (Inches(0.5),Inches(4.0)),(Inches(4.61),Inches(4.0)),(Inches(8.72),Inches(4.0))]
for i,(ic,ti,bo) in enumerate(feats):
    cx,cy=cols[i]; CARD(sl,ic,ti,bo,cx,cy,Inches(3.9),Inches(2.0))

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — DINE IN
# ════════════════════════════════════════════════════════════════════════════
sl = S(); HDR(sl,"Dine In","ขั้นตอนลูกค้าทานที่ร้าน")
flow=[("📷","สแกน QR"),("🪑","เลือกโต๊ะ"),("🍽️","เลือกเมนู"),
      ("🛒","ตรวจตะกร้า"),("💳","จ่ายเงิน"),("💬","เข้า LINE!")]
fw=Inches(1.85); fy=Inches(1.7)
for i,(ic,lb) in enumerate(flow):
    fx=Inches(0.5)+i*(fw+Inches(0.22))
    fc=RGBColor(0xEE,0xFF,0xDD) if i==5 else SAGE_L
    b=BOX(sl,fx,fy,fw,Inches(0.85),fill=fc,border=RGBColor(0xD4,0xE2,0xD2),bw=Pt(1))
    tf=b.text_frame; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=ic+"  "+lb; r.font.size=Pt(15); r.font.bold=True; r.font.color.rgb=DARK
    if i<5:
        T(sl,"→",fx+fw,fy+Inches(0.22),Inches(0.24),Inches(0.42),sz=Pt(18),c=SAGE_D,a=PP_ALIGN.CENTER)
steps_di=[
    ("1","สแกน QR ที่โต๊ะ — ไม่ต้องโหลดแอป","ใช้กล้องมือถือปกติได้เลย ทั้ง iPhone และ Android"),
    ("2","เลือกหมายเลขโต๊ะ","หมายเลขโต๊ะติดมากับออเดอร์อัตโนมัติ ร้านไม่ต้องถามซ้ำ"),
    ("3","เลือกเมนู → ตัวเลือก → ใส่ตะกร้า","ขนาด / ประเภทเนื้อ / ความหวาน — ราคาปรับทุกตัวเลือก"),
    ("4","ตรวจตะกร้า → ยืนยัน → เลือกวิธีชำระ","PromptPay หรือเงินสด → ออเดอร์เข้า LINE ร้านทันที"),
]
for i,(n,t,b) in enumerate(steps_di):
    STEP(sl,n,t,b,Inches(2.75)+i*Inches(1.12))

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — LINE ORDER (real data from actual order)
# ════════════════════════════════════════════════════════════════════════════
sl = S(); HDR(sl,"LINE OA","ออเดอร์จริงที่เข้ามาใน LINE ร้าน")

# Phone frame (left)
pf=BOX(sl,Inches(0.5),Inches(1.7),Inches(5.5),Inches(5.5),
       fill=WHITE,border=RGBColor(0x2A,0x2A,0x2A),bw=Pt(3))

# LINE header bar (dark)
lh=BOX(sl,Inches(0.5),Inches(1.7),Inches(5.5),Inches(0.65),fill=LINE_HD)
lh.line.fill.background()
T(sl,"←  ขาหมูนาย ต.  ☏  ⋮",Inches(0.55),Inches(1.77),Inches(5.4),Inches(0.5),
  sz=Pt(14),bold=True,c=WHITE)

# Chat bubble background
cb=BOX(sl,Inches(0.55),Inches(2.42),Inches(5.4),Inches(4.7),fill=RGBColor(0xF0,0xF0,0xF0))
cb.line.fill.background()

# Green bubble with real order data
bubble=BOX(sl,Inches(1.0),Inches(2.55),Inches(4.7),Inches(4.35),
           fill=LINE_GN,border=RGBColor(0x86,0xEF,0xAC),bw=Pt(1))
tf_b=bubble.text_frame; tf_b.word_wrap=True
order_lines=[
    ("📩 ระบบเมนู QR — ออเดอร์ใหม่", Pt(11), False, MID),
    ("🍖 ออเดอร์ใหม่!  โต๊ะ 6", Pt(16), True, DARK),
    ("ชุดไส้พะโล้เตาถ่าน  ·  ไส้ล้วน  ×2  →  200฿", Pt(12), False, DARK),
    ("ชุดขาหมูพะโล้สับ  ·  เนื้อหนัง  ×2  →  200฿", Pt(12), False, DARK),
    ("ชุดคากิพะโล้สับ  ×2  →  200฿", Pt(12), False, DARK),
    ("ไข่เปิ๊ดต้ม  ×7  →  70฿", Pt(12), False, DARK),
    ("ข้าวเปล่า  ×7  →  70฿", Pt(12), False, DARK),
    ("────────────────────────", Pt(9), False, RGBColor(0x86,0xEF,0xAC)),
    ("💰 ยอดรวม  740฿", Pt(16), True, DARK),
    ("💳 ชำระ: เงินสด  /  หมายเหตุ: เนื้อสัตว์กะเอิน", Pt(11), False, MID),
]
MT(sl,order_lines,Inches(1.05),Inches(2.6),Inches(4.55),Inches(4.2))

T(sl,"12:19 น.",Inches(4.8),Inches(6.82),Inches(0.8),Inches(0.3),sz=Pt(10),c=MID)

# Right panel tips
TIP(sl,"✅ ข้อมูลครบทุกอย่าง — ไม่ต้องถามลูกค้าเพิ่ม",
    "โต๊ะ · รายการ · ตัวเลือก · จำนวน · ยอดรวม · วิธีชำระ · หมายเหตุ\nร้านรับออเดอร์แล้วเริ่มทำได้เลยทันที",
    Inches(6.3),Inches(1.7),Inches(6.5),Inches(1.55),
    bg=GBG,tc=GGOOD,bc=GBORD)
TIP(sl,"💳 ลูกค้าเลือก QR PromptPay",
    "แสดง QR บนหน้าจอลูกค้า → สแกนโอน → แคปสลิปส่ง LINE OA\nร้านยืนยันรับเงินก่อนเริ่มทำ",
    Inches(6.3),Inches(3.45),Inches(6.5),Inches(1.55))
TIP(sl,"💬 LINE OA รับออเดอร์ได้ทันที",
    "ออเดอร์ทุกโต๊ะเข้า LINE เดียวกัน เรียงตามเวลา\nสามารถตรวจสอบย้อนหลังได้ตลอดเวลา",
    Inches(6.3),Inches(5.2),Inches(6.5),Inches(1.4),
    bg=HINT_BG,tc=HINT_FG,bc=HINT_BD)

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — DELIVERY
# ════════════════════════════════════════════════════════════════════════════
sl = S(); HDR(sl,"Delivery","ขั้นตอนลูกค้าสั่งส่งถึงบ้าน")
steps5=[
    ("1","กดปุ่ม 'ส่งถึงบ้าน' ที่หน้าแรก","ปุ่มอยู่ก่อนเลือกโต๊ะ — ลูกค้าสั่งจากบ้านได้เลย"),
    ("2","เลือกเมนูและตัวเลือก เหมือน Dine In","ใส่ตะกร้าได้หลายรายการ ราคาปรับตามตัวเลือก"),
    ("3","กรอกข้อมูลจัดส่ง","ชื่อ · เบอร์โทร · ที่อยู่ · เลือกรอบส่ง (เช้า 11:00 / บ่าย 14:00)"),
    ("4","ชำระ PromptPay → แคปสลิป → ส่ง LINE OA","พร้อมส่ง Pin โลเคชั่นด้วย เพื่อความสะดวก"),
    ("5","ร้านได้รับออเดอร์ + ที่อยู่ + สลิป ครบใน LINE","ยืนยันออเดอร์ → จัดส่งตามรอบที่ลูกค้าเลือก"),
]
for i,(n,t,b) in enumerate(steps5):
    STEP(sl,n,t,b,Inches(1.7)+i*Inches(1.0))
TIP(sl,"💡 แนะนำตั้ง 2 รอบส่ง: เช้า 11:00 + บ่าย 14:00",
    "ลูกค้าเลือกเอง ร้านรวมออเดอร์ส่งครั้งเดียว ประหยัดเวลาและค่าน้ำมัน",
    Inches(0.5),Inches(6.7),Inches(12.33),Inches(0.65),
    bg=HINT_BG,tc=HINT_FG,bc=HINT_BD)

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — CHAPTER ซ่อนเมนู
# ════════════════════════════════════════════════════════════════════════════
sl = S(); BG(sl,SAGE)
T(sl,"⭐  สำคัญที่สุด — เน้นเป็นพิเศษ",
  Inches(2),Inches(1.8),Inches(9.33),Inches(0.55),sz=Pt(15),bold=True,c=WHITE,a=PP_ALIGN.CENTER)
T(sl,"ซ่อน / แสดง\nเมนูชั่วคราว",
  Inches(1),Inches(2.45),Inches(11.33),Inches(2.1),sz=Pt(60),bold=True,c=WHITE,a=PP_ALIGN.CENTER)
T(sl,"เมื่อวัตถุดิบหมด หรืออยากปิดเมนูชั่วคราว\nทำได้เองใน 3 คลิก โดยไม่ต้องลบออกถาวร",
  Inches(2),Inches(4.85),Inches(9.33),Inches(1.0),sz=Pt(18),c=WHITE,a=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — ใช้ซ่อนเมื่อไร?
# ════════════════════════════════════════════════════════════════════════════
sl = S(); HDR(sl,"ซ่อนเมนู","ใช้ซ่อนเมนูเมื่อไร?")
scenarios=[
    ("🥩","วัตถุดิบหมดวันนั้น",
     "เช่น ขาหมูขายหมดแล้วก่อนร้านปิด → ซ่อนออกทันที ไม่ให้ลูกค้าสั่งเพิ่ม\nวันรุ่งขึ้นเปิดกลับมาได้ 1 คลิก"),
    ("⏸️","ปิดเมนูชั่วคราว",
     "เมนูฤดูกาล หรือรับออเดอร์ได้บางวัน\nซ่อนแทนการลบถาวร เปิดกลับมาได้ตลอด"),
    ("🔧","อยู่ระหว่างปรับราคา",
     "ซ่อนไว้ก่อนระหว่างแก้ราคา\nแก้เสร็จแล้วค่อยเปิดให้ลูกค้าเห็น"),
]
for i,(ic,ti,bo) in enumerate(scenarios):
    sy=Inches(1.75)+i*Inches(1.65)
    b=BOX(sl,Inches(0.5),sy,Inches(12.33),Inches(1.45),fill=SAGE_L); b.line.fill.background()
    MT(sl,[(ic+"  "+ti,Pt(19),True,DARK),(bo,Pt(14),False,MID)],
       Inches(0.65),sy+Inches(0.12),Inches(12.1),Inches(1.2))
TIP(sl,"✅ ซ่อน ≠ ลบ — ข้อมูลยังอยู่ครบทุกอย่าง",
    "รูปภาพ ราคา ตัวเลือก ทุกอย่างยังอยู่ในระบบ เพียงแต่ลูกค้าไม่เห็นชั่วคราว",
    Inches(0.5),Inches(6.7),Inches(12.33),Inches(0.65),bg=GBG,tc=GGOOD,bc=GBORD)

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — เปิด Supabase
# ════════════════════════════════════════════════════════════════════════════
sl = S(); HDR(sl,"ซ่อนเมนู · ขั้น 1","เปิดเว็บ Supabase")
steps8=[
    ("1","เปิดเบราว์เซอร์ → พิมพ์  supabase.com",""),
    ("2","กด Sign in → กรอก Email + Password ที่เราส่งให้ทาง LINE",""),
    ("3","คลิกชื่อโปรเจกต์ร้านของคุณ (เห็นในหน้าแรก)",""),
    ("4","เมนูซ้ายมือ → กด 'Table Editor' → คลิกตาราง  menu",""),
]
for i,(n,t,b) in enumerate(steps8):
    STEP(sl,n,t,b,Inches(1.75)+i*Inches(1.12))
TIP(sl,"💡 เคล็ดลับ: Bookmark URL ไว้เลย!",
    "เข้า Table Editor ครั้งแรก → กด Bookmark ในเบราว์เซอร์ทันที\nครั้งต่อไปกดบุ๊กมาร์กได้เลย ไม่ต้องคลิกหลายขั้น ประหยัดเวลามาก",
    Inches(0.5),Inches(6.35),Inches(12.33),Inches(0.9),bg=HINT_BG,tc=HINT_FG,bc=HINT_BD)

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — Table Editor หน้าตาจริง
# ════════════════════════════════════════════════════════════════════════════
sl = S(); HDR(sl,"ซ่อนเมนู · ขั้น 2","หน้าตา Table Editor ใน Supabase")
BG(sl,CREAM)
T(sl,"หลังกด Table Editor แล้วคลิก 'menu' จะเห็นตารางแบบนี้:",
  Inches(0.5),Inches(1.65),Inches(12.33),Inches(0.4),sz=Pt(15),c=MID)

# DB mockup
def db_bar_row(sl, y, h, bg_c, parts):
    b=BOX(sl,Inches(0.5),y,Inches(12.33),h,fill=bg_c); b.line.fill.background()
    tb=sl.shapes.add_textbox(Inches(0.6),y,Inches(12.1),h)
    tf=tb.text_frame; tf.word_wrap=False; p=tf.paragraphs[0]
    for t,c in parts:
        r=p.add_run(); r.text=t; r.font.size=Pt(13); r.font.color.rgb=c; r.font.name="Courier New"

db_bar_row(sl,Inches(2.1),Inches(0.55),DB_BG,[
    ("⚡ Supabase ",DB_GRN),("› ร้านขาหมูนาย ต. › Table Editor › ",DB_MUT),("menu",DB_GRN)])
db_bar_row(sl,Inches(2.65),Inches(0.4),DB_MID,[
    ("Authentication   ",DB_MUT),("● Table Editor   ",DB_GRN),("SQL Editor   Storage",DB_MUT)])
db_bar_row(sl,Inches(3.05),Inches(0.4),DB_MID,[
    ("id        ",DB_MUT),("name (ชื่อเมนู)                      ",DB_MUT),("price     ",DB_MUT),("is_available",DB_MUT)])
rows9=[("1 ","ขาหมูพะโล้เตาถ่าน (ชุดใหญ่)","180","✓ true",True),
       ("5 ","ชุดไส้พะโล้เตาถ่าน             ","100","✓ true",True),
       ("4 ","ข้าวขาหมูพะโล้เตาถ่าน          ","60 ","✓ true",True),
       ("7 ","กาแฟสด Premium                  ","65 ","✓ true",True)]
for ri,(id_,nm,pr,av,tval) in enumerate(rows9):
    ry=Inches(3.45)+ri*Inches(0.6)
    db_bar_row(sl,ry,Inches(0.58),DB_BG,[
        (id_.ljust(6),DB_TXT),(nm,DB_TXT),(pr.ljust(8),DB_TXT),(av,DB_GRN if tval else DB_RED)])

T(sl,"☝️  คอลัมน์ is_available ด้านขวา คือตัวควบคุม   true = แสดง   false = ซ่อน",
  Inches(0.5),Inches(5.87),Inches(12.33),Inches(0.5),sz=Pt(16),bold=True,c=RED)
TIP(sl,"ℹ️  is_available แปลว่า 'มีพร้อมสั่งไหม?'",
    "true = มีพร้อมสั่ง (ลูกค้าเห็น)  ·  false = ไม่มี (ลูกค้าไม่เห็น แต่ข้อมูลยังอยู่ในระบบ)",
    Inches(0.5),Inches(6.52),Inches(12.33),Inches(0.8))

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — คลิกเปลี่ยนค่า
# ════════════════════════════════════════════════════════════════════════════
sl = S(); HDR(sl,"ซ่อนเมนู · ขั้น 3","คลิกเซลล์ → พิมพ์ false → กด Enter")
BG(sl,CREAM)
T(sl,"ตัวอย่าง: ต้องการซ่อน 'ชุดไส้พะโล้เตาถ่าน' เพราะขายหมด",
  Inches(0.5),Inches(1.65),Inches(12.33),Inches(0.4),sz=Pt(15),c=MID)
db_bar_row(sl,Inches(2.1),Inches(0.5),DB_BG,[
    ("⚡ Supabase ",DB_GRN),("› Table Editor › menu",DB_MUT)])
db_bar_row(sl,Inches(2.6),Inches(0.4),DB_MID,[
    ("id        ",DB_MUT),("name (ชื่อเมนู)                      ",DB_MUT),("price     ",DB_MUT),("is_available",DB_MUT)])
rows10=[("1 ","ขาหมูพะโล้เตาถ่าน (ชุดใหญ่)","180","✓ true",True,False),
        ("5 ","ชุดไส้พะโล้เตาถ่าน             ","100","false|",False,True),
        ("4 ","ข้าวขาหมูพะโล้เตาถ่าน          ","60 ","✓ true",True,False)]
for ri,(id_,nm,pr,av,tval,hl) in enumerate(rows10):
    ry=Inches(3.0)+ri*Inches(0.65)
    rbg=RGBColor(0x0D,0x21,0x16) if hl else DB_BG
    b=BOX(sl,Inches(0.5),ry,Inches(12.33),Inches(0.62),fill=rbg)
    if hl: b.line.color.rgb=DB_GRN; b.line.width=Pt(2)
    else: b.line.fill.background()
    tb=sl.shapes.add_textbox(Inches(0.6),ry,Inches(12.1),Inches(0.6))
    tf=tb.text_frame; tf.word_wrap=False; p=tf.paragraphs[0]
    for t,c in [(id_.ljust(6),DB_TXT),(nm,CODE_YL if hl else DB_TXT),(pr.ljust(8),DB_TXT),(av,DB_RED if not tval else DB_GRN)]:
        r=p.add_run(); r.text=t; r.font.size=Pt(14); r.font.color.rgb=c
        r.font.bold=(hl and t.strip() in ["false|"]); r.font.name="Courier New"

T(sl,"⬆️  แถวนี้กำลังแก้ไข — พิมพ์  false  แล้วกด  Enter  เมนูหายทันที",
  Inches(0.5),Inches(4.98),Inches(12.33),Inches(0.5),sz=Pt(16),bold=True,c=RED)
STEP(sl,"①","ดับเบิลคลิกที่เซลล์ is_available ของเมนูที่ต้องการซ่อน",
     "เซลล์เข้าโหมดแก้ไข (เห็น cursor กะพริบ)",Inches(5.62))
STEP(sl,"②","ลบค่าเดิม พิมพ์ false (ตัวเล็กทั้งหมด) → กด Enter",
     "เมนูหายจากแอปลูกค้าทันที ไม่ต้อง Refresh หรือรอ",Inches(6.7))

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 11 — Before / After
# ════════════════════════════════════════════════════════════════════════════
sl = S(); HDR(sl,"ผลลัพธ์","ลูกค้าเห็นอะไร ก่อน / หลัง ซ่อนเมนู")

# Before
BOX(sl,Inches(0.5),Inches(1.7),Inches(5.9),Inches(4.0),
    fill=WHITE,border=RGBColor(0xFE,0xCA,0xCA))
BOX(sl,Inches(0.5),Inches(1.7),Inches(5.9),Inches(0.55),fill=RGBColor(0xFE,0xE2,0xE2))
T(sl,"ก่อนซ่อน ❌",Inches(0.5),Inches(1.7),Inches(5.9),Inches(0.55),
  sz=Pt(16),bold=True,c=RGBColor(0x99,0x1B,0x1B),a=PP_ALIGN.CENTER)
bf=[("🍖 ขาหมูพะโล้","180฿",False),("🫀 ชุดไส้พะโล้","100฿",True),
    ("🍚 ข้าวขาหมู","60฿",False),("☕ กาแฟสด","65฿",False)]
for bi,(nm,pr,hl) in enumerate(bf):
    by=Inches(2.35)+bi*Inches(0.72)
    if hl: BOX(sl,Inches(0.55),by,Inches(5.8),Inches(0.65),fill=RGBColor(0xFE,0xF9,0xC3))
    T(sl,nm,Inches(0.7),by+Inches(0.12),Inches(3.8),Inches(0.45),sz=Pt(16),bold=hl)
    T(sl,pr,Inches(5.4),by+Inches(0.12),Inches(0.8),Inches(0.45),sz=Pt(16),bold=hl,a=PP_ALIGN.RIGHT)

# After
BOX(sl,Inches(6.7),Inches(1.7),Inches(6.13),Inches(4.0),
    fill=WHITE,border=RGBColor(0x86,0xEF,0xAC))
BOX(sl,Inches(6.7),Inches(1.7),Inches(6.13),Inches(0.55),fill=RGBColor(0xDC,0xFC,0xE7))
T(sl,"หลังซ่อน ✅",Inches(6.7),Inches(1.7),Inches(6.13),Inches(0.55),
  sz=Pt(16),bold=True,c=GGOOD,a=PP_ALIGN.CENTER)
af=[("🍖 ขาหมูพะโล้","180฿"),("🍚 ข้าวขาหมู","60฿"),("☕ กาแฟสด","65฿")]
for ai,(nm,pr) in enumerate(af):
    ay=Inches(2.35)+ai*Inches(0.72)
    T(sl,nm,Inches(6.85),ay+Inches(0.12),Inches(4.0),Inches(0.45),sz=Pt(16))
    T(sl,pr,Inches(11.65),ay+Inches(0.12),Inches(0.9),Inches(0.45),sz=Pt(16),a=PP_ALIGN.RIGHT)
T(sl,"✓  'ชุดไส้' หายออกไปแล้ว",Inches(6.85),Inches(4.56),Inches(5.8),Inches(0.4),
  sz=Pt(14),bold=True,c=RGBColor(0x16,0xA3,0x4A))

TIP(sl,"↩️ อยากเปิดกลับมา?",
    "กลับ Supabase → ดับเบิลคลิกเซลล์เดิม → พิมพ์  true → Enter  เมนูกลับมาทันที",
    Inches(0.5),Inches(5.88),Inches(12.33),Inches(0.78),bg=WARN_BG,tc=WARN_FG,bc=WARN_BD)
TIP(sl,"✅ ซ่อนไม่ลบ — ข้อมูลไม่หายไปไหน",
    "รูปภาพ ราคา ตัวเลือก ทุกอย่างยังอยู่ครบ แค่ลูกค้าไม่เห็นชั่วคราว",
    Inches(0.5),Inches(6.75),Inches(12.33),Inches(0.62),bg=GBG,tc=GGOOD,bc=GBORD)

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 12 — SQL เร็วกว่า
# ════════════════════════════════════════════════════════════════════════════
sl = S(); HDR(sl,"💡 ข้อเสนอแนะ","วิธีเร็วกว่า — ใช้ SQL แทน Table Editor")
TIP(sl,"🎯 ทำไม SQL ถึงง่ายกว่า?",
    "Table Editor: ต้องหาแถว → เลื่อนหา column → คลิก → แก้\nSQL: พิมพ์ชื่อเมนู → กด Run → เสร็จ ใช้เวลาไม่ถึง 10 วินาที!",
    Inches(0.5),Inches(1.7),Inches(12.33),Inches(1.2),bg=HINT_BG,tc=HINT_FG,bc=HINT_BD)
T(sl,"เมนูซ้าย → SQL Editor → New Query → วางโค้ดนี้ → กด Run",
  Inches(0.5),Inches(3.05),Inches(12.33),Inches(0.42),sz=Pt(15),c=MID)

sql=BOX(sl,Inches(0.5),Inches(3.55),Inches(12.33),Inches(2.65),fill=CODE_BG)
sql.line.fill.background()
sql_lines=[
    [("-- 🙈 ซ่อนเมนู (เปลี่ยนแค่ชื่อในกรอบเหลือง)",CODE_GR,False)],
    [("UPDATE ",CODE_PR,True),("public.menu",CODE_BL,False)],
    [("SET ",CODE_PR,True),("is_available = ",CODE_BL,False),("false",CODE_GN,True)],
    [("WHERE ",CODE_PR,True),("name ",CODE_BL,False),("ILIKE ",CODE_PR,True),("'%",CODE_GN,False),("ชุดไส้",CODE_YL,True),("%';",CODE_GN,False)],
    [("",CODE_GR,False)],
    [("-- 👁️ เปิดกลับมา",CODE_GR,False)],
    [("UPDATE ",CODE_PR,True),("public.menu  SET is_available = ",CODE_BL,False),("true  ",CODE_GN,True),("WHERE name ILIKE ",CODE_PR,False),("'%",CODE_GN,False),("ชุดไส้",CODE_YL,True),("%';",CODE_GN,False)],
]
tb_sql=sl.shapes.add_textbox(Inches(0.65),Inches(3.65),Inches(12.0),Inches(2.4))
tf_sql=tb_sql.text_frame; tf_sql.word_wrap=False; first=True
for parts in sql_lines:
    p=tf_sql.paragraphs[0] if first else tf_sql.add_paragraph(); first=False
    for t,c,b in parts:
        r=p.add_run(); r.text=t; r.font.size=Pt(15); r.font.color.rgb=c
        r.font.bold=b; r.font.name="Courier New"

TIP(sl,"✏️  เปลี่ยนแค่คำสีเหลือง  ชุดไส้  = ชื่อ (หรือส่วนของชื่อ) เมนูที่ต้องการซ่อน",
    "เช่น ซ่อน 'ขาหมูพะโล้เตาถ่าน' → พิมพ์  ขาหมู   ILIKE ค้นหาชื่อที่มีคำนั้นอยู่ ไม่สนตัวเล็ก/ใหญ่",
    Inches(0.5),Inches(6.35),Inches(12.33),Inches(0.9),bg=HINT_BG,tc=HINT_FG,bc=HINT_BD)

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 13 — แก้ราคา
# ════════════════════════════════════════════════════════════════════════════
sl = S(); HDR(sl,"แก้ราคา","แก้ราคาเมนูด้วยตัวเอง")
steps13=[
    ("1","Supabase → Table Editor → ตาราง menu","เห็นรายการเมนูทั้งหมด พร้อมราคาปัจจุบัน"),
    ("2","หาเมนูที่ต้องการ → ดับเบิลคลิกเซลล์ price","เซลล์จะเข้าโหมดแก้ไข เห็น cursor กะพริบ"),
    ("3","พิมพ์ราคาใหม่ → กด Enter","แอปลูกค้าอัปเดตภายใน 1–2 วินาที ไม่ต้อง deploy ใหม่"),
]
for i,(n,t,b) in enumerate(steps13):
    STEP(sl,n,t,b,Inches(1.75)+i*Inches(1.1))
db_bar_row(sl,Inches(5.0),Inches(0.5),DB_BG,[
    ("⚡ Supabase ",DB_GRN),("› Table Editor › menu",DB_MUT)])
db_bar_row(sl,Inches(5.5),Inches(0.4),DB_MID,[
    ("id      ",DB_MUT),("name                                  ",DB_MUT),("price        ",DB_MUT),("is_available",DB_MUT)])
rows13=[("1 ","ขาหมูพะโล้เตาถ่าน","200|",True),("4 ","ข้าวขาหมูพะโล้","60 ",False)]
for ri,(id_,nm,pr,hl) in enumerate(rows13):
    ry=Inches(5.9)+ri*Inches(0.62)
    rbg=RGBColor(0x0D,0x21,0x16) if hl else DB_BG
    b=BOX(sl,Inches(0.5),ry,Inches(12.33),Inches(0.6),fill=rbg)
    if hl: b.line.color.rgb=DB_GRN; b.line.width=Pt(2)
    else: b.line.fill.background()
    tb=sl.shapes.add_textbox(Inches(0.6),ry,Inches(12.1),Inches(0.58))
    tf=tb.text_frame; p=tf.paragraphs[0]
    for t,c in [(id_.ljust(6),DB_TXT),(nm.ljust(36),DB_TXT),(pr.ljust(14),CODE_YL if hl else DB_TXT),("✓ true",DB_GRN)]:
        r=p.add_run(); r.text=t; r.font.size=Pt(14); r.font.color.rgb=c; r.font.name="Courier New"
T(sl,"☝️  คลิกเซลล์ price → พิมพ์ราคาใหม่ (เช่น 200) → Enter เสร็จ อัปเดตทันที",
  Inches(0.5),Inches(7.18),Inches(12.33),Inches(0.22),sz=Pt(14),bold=True,c=RED)

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 14 — Supabase Paused
# ════════════════════════════════════════════════════════════════════════════
sl = S(); HDR(sl,"แก้ปัญหา","Supabase ขึ้น 'Paused' — ทำอย่างไร?")
TIP(sl,"⚠️  Paused เกิดจากอะไร?",
    "Supabase แบบฟรีจะหยุดทำงานอัตโนมัติถ้าไม่ได้ใช้งานนาน 7 วัน\nเมนูลูกค้ายังดูได้ แต่การแก้ราคาจะยังไม่มีผล จนกว่าจะ Resume",
    Inches(0.5),Inches(1.7),Inches(12.33),Inches(1.2),bg=WARN_BG,tc=WARN_FG,bc=WARN_BD)
steps14=[
    ("1","เข้า supabase.com → Sign in → คลิกโปรเจกต์","จะเห็นสถานะ 'Paused' หรือ 'Unhealthy'"),
    ("2","กดปุ่มสีเขียว 'Resume project'","โปรเจกต์จะเริ่มต้นใหม่ ใช้เวลา 1–2 นาที"),
    ("3","รอ 2 นาที → กด Refresh หน้าแอปเมนู","ทุกอย่างกลับมาปกติ ราคาอัปเดตได้ตามเดิม"),
]
for i,(n,t,b) in enumerate(steps14):
    STEP(sl,n,t,b,Inches(3.1)+i*Inches(1.1))
TIP(sl,"💡 ป้องกัน Paused — เข้า Supabase สัปดาห์ละครั้ง",
    "แค่เปิดหน้า Table Editor ก็นับว่า active แล้ว ไม่ต้องทำอะไรเพิ่ม",
    Inches(0.5),Inches(6.45),Inches(12.33),Inches(0.75),bg=HINT_BG,tc=HINT_FG,bc=HINT_BD)

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 15 — จบ
# ════════════════════════════════════════════════════════════════════════════
sl = S(); BG(sl,SAGE)
T(sl,"🎥",Inches(5.66),Inches(0.85),Inches(2),Inches(0.9),
  sz=Pt(52),c=WHITE,a=PP_ALIGN.CENTER)
T(sl,"จบสไลด์แล้ว!",Inches(2),Inches(1.8),Inches(9.33),Inches(0.9),
  sz=Pt(48),bold=True,c=WHITE,a=PP_ALIGN.CENTER)
T(sl,"ดูครบแล้ว ถ้ามีคำถามอะไรเพิ่ม ทักเราทาง LINE ได้เลยครับ",
  Inches(2),Inches(2.85),Inches(9.33),Inches(0.55),sz=Pt(18),c=WHITE,a=PP_ALIGN.CENTER)

zoom=BOX(sl,Inches(2.5),Inches(3.55),Inches(8.33),Inches(3.4),fill=WHITE)
zoom.line.color.rgb=WHITE
tb_z=sl.shapes.add_textbox(Inches(2.65),Inches(3.65),Inches(8.0),Inches(3.1))
tf_z=tb_z.text_frame; tf_z.word_wrap=True
zoom_lines=[
    ("📅 อย่าลืมนัด Zoom สอนงานจริง 30 นาที", Pt(17), True, DARK),
    ("✓  ทดลองซ่อนเมนูพร้อมกัน live", Pt(16), False, DARK),
    ("✓  ทดลองแก้ราคาใน Supabase", Pt(16), False, DARK),
    ("✓  ดูออเดอร์เข้า LINE OA จริง", Pt(16), False, DARK),
    ("✓  ตอบคำถามที่ค้างอยู่ทุกข้อ", Pt(16), False, DARK),
    ("✓  บันทึก Zoom ไว้ดูย้อนหลังได้", Pt(16), False, DARK),
]
first=True
for t,sz,b,c in zoom_lines:
    p=tf_z.paragraphs[0] if first else tf_z.add_paragraph(); first=False
    r=p.add_run(); r.text=t; r.font.size=sz; r.font.bold=b; r.font.color.rgb=c

prs.save("/home/user/khaomoo-menu/docs/slides-premium-v2.pptx")
print("PPTX saved.")
