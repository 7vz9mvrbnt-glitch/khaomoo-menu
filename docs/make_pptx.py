#!/usr/bin/env python3
"""Generate slides-premium.pptx from slide content."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── Helpers ─────────────────────────────────────────────────────────────────
W, H = Inches(10), Inches(7.5)

SAGE       = RGBColor(0x87, 0xA0, 0x80)
SAGE_L     = RGBColor(0xEC, 0xF2, 0xEB)
SAGE_D     = RGBColor(0x55, 0x72, 0x52)
CREAM      = RGBColor(0xF8, 0xF5, 0xEF)
DARK       = RGBColor(0x2A, 0x2A, 0x2A)
MID        = RGBColor(0x5E, 0x5E, 0x5E)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
RED        = RGBColor(0xDC, 0x26, 0x26)
GREEN_GOOD = RGBColor(0x16, 0x65, 0x34)
WARN_BG    = RGBColor(0xFF, 0xF9, 0xE6)
WARN_FG    = RGBColor(0x7A, 0x55, 0x00)
HINT_BG    = RGBColor(0xEE, 0xF2, 0xFF)
HINT_FG    = RGBColor(0x37, 0x30, 0xA3)
DB_BG      = RGBColor(0x0F, 0x17, 0x2A)
DB_MID     = RGBColor(0x1E, 0x29, 0x3B)
DB_GREEN   = RGBColor(0x3E, 0xCF, 0x8E)
DB_RED     = RGBColor(0xF8, 0x71, 0x71)
DB_TEXT    = RGBColor(0xCB, 0xD5, 0xE1)
DB_MUTED   = RGBColor(0x94, 0xA3, 0xB8)
CODE_BG    = RGBColor(0x1E, 0x29, 0x3B)
CODE_BLUE  = RGBColor(0x7D, 0xD3, 0xFC)
CODE_GREEN = RGBColor(0x86, 0xEF, 0xAC)
CODE_PURP  = RGBColor(0xC0, 0x84, 0xFC)
CODE_GRAY  = RGBColor(0x47, 0x55, 0x69)
CODE_YEL   = RGBColor(0xFC, 0xD3, 0x4D)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # blank


def slide():
    return prs.slides.add_slide(blank_layout)


def bg(sl, color):
    fill = sl.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(sl, x, y, w, h, fill=None, border=None, border_w=Pt(1)):
    shape = sl.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = border_w
    else:
        shape.line.fill.background()
    return shape


def txt(sl, text, x, y, w, h,
        size=Pt(14), bold=False, color=DARK, align=PP_ALIGN.LEFT,
        wrap=True, italic=False):
    txb = sl.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def mtxt(sl, lines, x, y, w, h, default_size=Pt(13), default_color=DARK, default_bold=False, align=PP_ALIGN.LEFT):
    """lines = list of (text, size, bold, color, italic)"""
    txb = sl.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            item = (item, default_size, default_bold, default_color, False)
        t, sz, b, c, it = item
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = t
        run.font.size = sz
        run.font.bold = b
        run.font.color.rgb = c
        run.font.italic = it
    return txb


def tag(sl, label, x, y):
    b = box(sl, x, y, Inches(1.6), Inches(0.28), fill=SAGE)
    tf = b.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(8)
    r.font.bold = True
    r.font.color.rgb = WHITE
    return b


def hline(sl, x, y, w, color=SAGE_L, height=Pt(1)):
    ln = box(sl, x, y, w, Inches(0.01), fill=color)
    return ln


def step_circle(sl, num_str, x, y):
    c = box(sl, x, y, Inches(0.38), Inches(0.38), fill=SAGE)
    tf = c.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = num_str
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = WHITE
    return c


def tip_box(sl, title, body, x, y, w, h, bg_color=SAGE_L, title_color=SAGE_D, border_color=SAGE):
    b = box(sl, x, y, w, h, fill=bg_color, border=border_color, border_w=Pt(2))
    # left accent bar
    box(sl, x, y, Inches(0.05), h, fill=border_color)
    tf = b.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT
    r0 = p0.add_run()
    r0.text = title
    r0.font.size = Pt(11)
    r0.font.bold = True
    r0.font.color.rgb = title_color
    p1 = tf.add_paragraph()
    r1 = p1.add_run()
    r1.text = body
    r1.font.size = Pt(10)
    r1.font.color.rgb = DARK
    return b


def card(sl, icon, title, body, x, y, w, h):
    b = box(sl, x, y, w, h, fill=SAGE_L)
    b.line.fill.background()
    tf = b.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]; r0 = p0.add_run()
    r0.text = icon + "  " + title
    r0.font.size = Pt(11); r0.font.bold = True; r0.font.color.rgb = DARK
    p1 = tf.add_paragraph(); r1 = p1.add_run()
    r1.text = body
    r1.font.size = Pt(9.5); r1.font.color.rgb = MID
    return b


def slide_header(sl, tag_label, title, bg_c=WHITE):
    bg(sl, bg_c)
    tag(sl, tag_label, Inches(0.4), Inches(0.32))
    txt(sl, title, Inches(0.4), Inches(0.62), Inches(9.2), Inches(0.55),
        size=Pt(22), bold=True, color=DARK)
    hline(sl, Inches(0.4), Inches(1.22), Inches(9.2), color=SAGE_L, height=Pt(2))


def step_row(sl, num, title, body, y):
    step_circle(sl, str(num), Inches(0.4), y)
    txt(sl, title, Inches(0.85), y - Inches(0.01), Inches(8.8), Inches(0.3),
        size=Pt(13), bold=True, color=DARK)
    txt(sl, body, Inches(0.85), y + Inches(0.28), Inches(8.8), Inches(0.4),
        size=Pt(10.5), color=MID)


# ════════════════════════════════════════════════════════════
#  SLIDE 1 — COVER
# ════════════════════════════════════════════════════════════
s = slide(); bg(s, SAGE)
txt(s, "สไลด์สอนงาน · แพ็ก Premium",
    Inches(2.5), Inches(1.2), Inches(5), Inches(0.4),
    size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "🚀", Inches(4.5), Inches(1.75), Inches(1), Inches(0.7),
    size=Pt(36), color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "สอนใช้งาน\nแพ็ก Premium",
    Inches(1.5), Inches(2.5), Inches(7), Inches(1.4),
    size=Pt(38), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "QR Code เมนูครบวงจร  ·  ออเดอร์เข้า LINE  ·  Delivery  ·  2 ภาษา",
    Inches(1), Inches(4.0), Inches(8), Inches(0.5),
    size=Pt(14), color=WHITE, align=PP_ALIGN.CENTER)
b = box(s, Inches(2.8), Inches(4.8), Inches(4.4), Inches(0.55),
        fill=RGBColor(0xFF,0xFF,0xFF), border=RGBColor(0xFF,0xFF,0xFF))
tf2 = b.text_frame; p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run(); r2.text = "กด  ←  →  เพื่อเปลี่ยนสไลด์"
r2.font.size = Pt(11); r2.font.color.rgb = SAGE_D

# ════════════════════════════════════════════════════════════
#  SLIDE 2 — ฟีเจอร์
# ════════════════════════════════════════════════════════════
s = slide(); slide_header(s, "ภาพรวม", "แพ็ก Premium — ได้อะไรบ้าง?")
feats = [
    ("💬","ออเดอร์เข้า LINE OA","ลูกค้ากดยืนยัน → ร้านได้รับออเดอร์ครบใน LINE ทันที ไม่มีจดผิด"),
    ("🏠","Delivery ส่งถึงบ้าน","กรอกที่อยู่ + ชำระ PromptPay ครบใน 1 ขั้นตอน"),
    ("🌐","2 ภาษา ไทย / EN","กดปุ่มสลับ ทุกหน้าเปลี่ยนเป็นอังกฤษพร้อมกัน"),
    ("🎛️","ตัวเลือกเมนูครบ","ขนาด · ประเภทเนื้อ · ท็อปปิ้ง · ความหวาน ราคาปรับอัตโนมัติ"),
    ("✏️","แก้ราคาเองได้","Supabase Dashboard — คลิก แก้ Enter เสร็จ ไม่ต้องรอเรา"),
    ("👁️","ซ่อน/แสดงเมนู","ของหมด — ซ่อนได้ 1 คลิก โดยไม่ต้องลบออกถาวร"),
]
cols = [(Inches(0.35), Inches(1.4)), (Inches(5.1), Inches(1.4)),
        (Inches(0.35), Inches(3.0)), (Inches(5.1), Inches(3.0)),
        (Inches(0.35), Inches(4.6)), (Inches(5.1), Inches(4.6))]
for i, (icon, title, body) in enumerate(feats):
    cx, cy = cols[i]
    card(s, icon, title, body, cx, cy, Inches(4.55), Inches(1.4))

# ════════════════════════════════════════════════════════════
#  SLIDE 3 — DINE IN FLOW
# ════════════════════════════════════════════════════════════
s = slide(); slide_header(s, "Dine In", "ขั้นตอนลูกค้าทานที่ร้าน")
flow_items = [("📷","สแกน QR"),("🪑","เลือกโต๊ะ"),("🍽️","เลือกเมนู"),
              ("🛒","ตรวจตะกร้า"),("💳","จ่ายเงิน"),("💬","เข้า LINE ร้าน!")]
fw = Inches(1.35); fgap = Inches(0.12); fy = Inches(1.4)
for i, (icon, label) in enumerate(flow_items):
    fx = Inches(0.35) + i*(fw + fgap + Inches(0.15))
    if i < 5:
        fc = SAGE_L
    else:
        fc = RGBColor(0xEE, 0xFF, 0xDD)
    b2 = box(s, fx, fy, fw, Inches(0.8), fill=fc, border=RGBColor(0xD4,0xE2,0xD2))
    tf3 = b2.text_frame; tf3.word_wrap = False
    p3 = tf3.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run(); r3.text = icon + "  " + label
    r3.font.size = Pt(11); r3.font.bold = True; r3.font.color.rgb = DARK
    if i < 5:
        txt(s, "→", fx + fw + Inches(0.01), fy + Inches(0.22), Inches(0.18), Inches(0.36),
            size=Pt(14), color=SAGE_D, align=PP_ALIGN.CENTER)

steps_di = [
    ("1","เปิดกล้องมือถือ → สแกน QR ที่โต๊ะ","ไม่ต้องโหลดแอป ใช้กล้องปกติได้เลย ทั้ง iPhone และ Android"),
    ("2","เลือกหมายเลขโต๊ะ","โต๊ะติดมากับออเดอร์อัตโนมัติ ร้านไม่ต้องถามใหม่"),
    ("3","เลือกเมนู → เลือกตัวเลือก → เพิ่มลงตะกร้า","ขนาด / ประเภทเนื้อ / ความหวาน — ราคาปรับอัตโนมัติทุกตัวเลือก"),
    ("4","ตรวจรายการ → กดยืนยัน → เลือกวิธีชำระ","QR PromptPay หรือเงินสด — ออเดอร์เข้า LINE ร้านทันที"),
]
for i, (n, t, b) in enumerate(steps_di):
    step_row(s, n, t, b, Inches(2.4) + i*Inches(1.2))

# ════════════════════════════════════════════════════════════
#  SLIDE 4 — LINE OA ออเดอร์
# ════════════════════════════════════════════════════════════
s = slide(); slide_header(s, "LINE OA", "ออเดอร์ที่เข้ามาใน LINE ร้าน")
# LINE bubble mockup
lb = box(s, Inches(0.4), Inches(1.45), Inches(5.2), Inches(4.5),
         fill=RGBColor(0xEE,0xFF,0xDD), border=RGBColor(0x86,0xEF,0xAC))
tf_lb = lb.text_frame; tf_lb.word_wrap = True
lines_lb = [
    ("📩 ออเดอร์ใหม่ · จากระบบเมนู QR", Pt(9), False, MID, False),
    ("🍖 ออเดอร์ใหม่! โต๊ะ 3", Pt(14), True, DARK, False),
    ("ขาหมูพะโล้เตาถ่าน (ชุดใหญ่)  .....  180฿", Pt(11), False, DARK, False),
    ("ข้าวขาหมูพะโล้ × 2  ..................  120฿", Pt(11), False, DARK, False),
    ("กาแฟลาเต้เย็น (M) หวานน้อย  ....  75฿", Pt(11), False, DARK, False),
    ("─────────────────────────────", Pt(8), False, RGBColor(0x86,0xEF,0xAC), False),
    ("💰 ยอดรวม  ..............................  375฿", Pt(13), True, DARK, False),
    ("💳 ชำระ  .......................  QR PromptPay", Pt(10), False, MID, False),
    ("📝 หมายเหตุ  ......................  ไม่ใส่ผัก", Pt(10), False, MID, False),
]
first_lb = True
for item in lines_lb:
    t2, sz2, b2, c2, it2 = item
    if first_lb:
        p2 = tf_lb.paragraphs[0]; first_lb = False
    else:
        p2 = tf_lb.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run(); r2.text = t2
    r2.font.size = sz2; r2.font.bold = b2; r2.font.color.rgb = c2

tip_box(s, "✅ ข้อมูลครบทุกอย่าง ไม่ต้องถามเพิ่ม",
        "โต๊ะ · รายการ · ราคา · วิธีชำระ · หมายเหตุ\nร้านรับออเดอร์และเริ่มทำได้เลย",
        Inches(5.8), Inches(1.45), Inches(3.9), Inches(1.4),
        bg_color=RGBColor(0xF0,0xFD,0xF4), title_color=GREEN_GOOD,
        border_color=RGBColor(0x22,0xC5,0x5E))
tip_box(s, "💳 ลูกค้าเลือก QR PromptPay",
        "แสดง QR บนหน้าจอ → ลูกค้าสแกนโอน\nแคปสลิปส่ง LINE OA → ร้านยืนยัน",
        Inches(5.8), Inches(3.0), Inches(3.9), Inches(1.4))

# ════════════════════════════════════════════════════════════
#  SLIDE 5 — DELIVERY
# ════════════════════════════════════════════════════════════
s = slide(); slide_header(s, "Delivery", "ขั้นตอนลูกค้าสั่งส่งถึงบ้าน")
steps_del = [
    ("1","หน้าแรก → กดปุ่ม 'ส่งถึงบ้าน'","ปุ่มอยู่หน้าแรก ก่อนเลือกโต๊ะ"),
    ("2","เลือกเมนูและตัวเลือก เหมือน Dine In","ใส่ตะกร้าได้หลายรายการ"),
    ("3","กรอกข้อมูลจัดส่ง","ชื่อ · เบอร์โทร · ที่อยู่ · เลือกรอบส่ง (เช้า 11:00 / บ่าย 14:00)"),
    ("4","ชำระ QR PromptPay → แคปสลิป → ส่ง LINE OA ร้าน","พร้อมส่ง Pin โลเคชั่นบ้านด้วย"),
    ("5","ร้านได้รับออเดอร์ + ที่อยู่ + สลิป ครบใน LINE","ยืนยันออเดอร์ → จัดส่งตามรอบที่ลูกค้าเลือก"),
]
for i, (n, t, b) in enumerate(steps_del):
    step_row(s, n, t, b, Inches(1.4) + i*Inches(1.0))
tip_box(s, "💡 แนะนำตั้ง 2 รอบส่ง",
        "เช้า (11:00) และบ่าย (14:00) ลูกค้าเลือกเอง\nร้านรวมออเดอร์ส่งครั้งเดียว ประหยัดเวลา",
        Inches(0.4), Inches(6.55), Inches(9.2), Inches(0.8),
        bg_color=HINT_BG, title_color=HINT_FG,
        border_color=RGBColor(0x63,0x66,0xF1))

# ════════════════════════════════════════════════════════════
#  SLIDE 6 — CHAPTER ซ่อนเมนู
# ════════════════════════════════════════════════════════════
s = slide(); bg(s, SAGE)
txt(s, "⭐ สำคัญ — เน้นเป็นพิเศษ",
    Inches(1), Inches(2.0), Inches(8), Inches(0.4),
    size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "ซ่อน / แสดง\nเมนูชั่วคราว",
    Inches(1), Inches(2.5), Inches(8), Inches(1.8),
    size=Pt(40), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "เมื่อวัตถุดิบหมด หรืออยากปิดเมนูชั่วคราว\nทำได้เองใน 3 คลิก โดยไม่ต้องลบ",
    Inches(1.5), Inches(4.45), Inches(7), Inches(0.8),
    size=Pt(14), color=WHITE, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
#  SLIDE 7 — ใช้ซ่อนเมนูเมื่อไร?
# ════════════════════════════════════════════════════════════
s = slide(); slide_header(s, "ซ่อนเมนู", "ใช้ซ่อนเมนูเมื่อไร?")
scenarios = [
    ("🥩","วัตถุดิบหมดวันนั้น","เช่น ขาหมูขายหมดแล้วก่อนร้านปิด ซ่อนออกไม่ให้ลูกค้าสั่งเพิ่ม วันรุ่งขึ้นเปิดกลับมา"),
    ("⏸️","ปิดเมนูชั่วคราว","เมนูฤดูกาล หรือรับออเดอร์ได้บางวัน ซ่อนแทนการลบ เปิดได้ตลอดเวลา"),
    ("🔧","อยู่ระหว่างปรับราคา","ซ่อนไว้ก่อนระหว่างแก้ราคา แก้เสร็จแล้วค่อยเปิดให้ลูกค้าเห็น"),
]
for i, (icon, title, body) in enumerate(scenarios):
    sy = Inches(1.4) + i*Inches(1.55)
    sb2 = box(s, Inches(0.35), sy, Inches(9.3), Inches(1.35), fill=SAGE_L)
    sb2.line.fill.background()
    tf_s = sb2.text_frame; tf_s.word_wrap = True
    ps0 = tf_s.paragraphs[0]; ps0.alignment = PP_ALIGN.LEFT
    rs0 = ps0.add_run()
    rs0.text = icon + "  " + title
    rs0.font.size = Pt(13); rs0.font.bold = True; rs0.font.color.rgb = DARK
    ps1 = tf_s.add_paragraph()
    rs1 = ps1.add_run(); rs1.text = body
    rs1.font.size = Pt(11); rs1.font.color.rgb = MID
tip_box(s, "✅ ซ่อน ≠ ลบ — ข้อมูลยังอยู่ครบ",
        "เปลี่ยนกลับได้ตลอดเวลา ลูกค้าไม่เห็น แต่ข้อมูลในระบบยังอยู่ครบ",
        Inches(0.35), Inches(6.15), Inches(9.3), Inches(0.9),
        bg_color=RGBColor(0xF0,0xFD,0xF4), title_color=GREEN_GOOD,
        border_color=RGBColor(0x22,0xC5,0x5E))

# ════════════════════════════════════════════════════════════
#  SLIDE 8 — เปิด Supabase
# ════════════════════════════════════════════════════════════
s = slide(); slide_header(s, "ซ่อนเมนู · ขั้น 1", "เปิดเว็บ Supabase")
steps_sb = [
    ("1","เปิดเบราว์เซอร์ → พิมพ์ supabase.com",""),
    ("2","กด Sign in → กรอก Email + Password","ใช้ email และ password ที่เราส่งให้ทาง LINE"),
    ("3","คลิกชื่อโปรเจกต์ร้านของคุณ","จะเห็นชื่อร้านในหน้าแรก กดคลิกเข้าไป"),
    ("4","เมนูซ้ายมือ → กด 'Table Editor' → คลิกตาราง menu","จะเห็นรายการตาราง → คลิกตาราง menu"),
]
for i, (n, t, b) in enumerate(steps_sb):
    step_row(s, n, t, b, Inches(1.4) + i*Inches(1.2))
tip_box(s, "💡 เคล็ดลับ: Bookmark ไว้เลย!",
        "เข้า Table Editor ครั้งแรกแล้ว กด Bookmark ในเบราว์เซอร์ไว้เลย\nครั้งต่อไปกดบุ๊กมาร์กได้ทันที ไม่ต้องคลิกหลายขั้น",
        Inches(0.4), Inches(6.3), Inches(9.2), Inches(1.0),
        bg_color=HINT_BG, title_color=HINT_FG,
        border_color=RGBColor(0x63,0x66,0xF1))

# ════════════════════════════════════════════════════════════
#  SLIDE 9 — Table Editor หน้าตาจริง
# ════════════════════════════════════════════════════════════
s = slide(); slide_header(s, "ซ่อนเมนู · ขั้น 2", "หน้าตา Table Editor จริง ๆ")
bg(s, CREAM)
# DB header bar
db_bar = box(s, Inches(0.35), Inches(1.4), Inches(9.3), Inches(0.45), fill=DB_BG)
db_bar.line.fill.background()
tf_db = db_bar.text_frame; tf_db.word_wrap = False
p_db = tf_db.paragraphs[0]; p_db.alignment = PP_ALIGN.LEFT
for txt_part, color in [("⚡ Supabase", DB_GREEN), ("  ›  ร้านข้าวขาหมู  ›  Table Editor  ›  ", DB_MUTED), ("menu", DB_GREEN)]:
    r_db = p_db.add_run(); r_db.text = txt_part
    r_db.font.size = Pt(10); r_db.font.bold = (color == DB_GREEN); r_db.font.color.rgb = color

# nav bar
nav_bar = box(s, Inches(0.35), Inches(1.85), Inches(9.3), Inches(0.35), fill=DB_MID)
nav_bar.line.fill.background()
tf_nav = nav_bar.text_frame; tf_nav.word_wrap = False
p_nav = tf_nav.paragraphs[0]; p_nav.alignment = PP_ALIGN.LEFT
for txt_part, color in [("Authentication   ", DB_MUTED), ("● Table Editor   ", DB_GREEN), ("SQL Editor   Storage", DB_MUTED)]:
    r_nav = p_nav.add_run(); r_nav.text = txt_part
    r_nav.font.size = Pt(9); r_nav.font.color.rgb = color

# thead
thead = box(s, Inches(0.35), Inches(2.2), Inches(9.3), Inches(0.35), fill=DB_MID)
thead.line.fill.background()
tf_th = thead.text_frame; tf_th.word_wrap = False
p_th = tf_th.paragraphs[0]
cols_w = [0.6, 3.2, 1.0, 2.0]  # inches
for col_name in ["id", "name (ชื่อเมนู)", "price", "is_available"]:
    r_th = p_th.add_run(); r_th.text = col_name.ljust(22)
    r_th.font.size = Pt(9); r_th.font.color.rgb = DB_MUTED; r_th.font.bold = True

# data rows
rows_data = [
    ("1","ขาหมูพะโล้เตาถ่าน","180","✓ true", True, False),
    ("5","ชุดไส้พะโล้เตาถ่าน","120","✓ true", True, False),
    ("4","ข้าวขาหมูพะโล้","60","✓ true", True, False),
    ("7","กาแฟสด Premium","65","✓ true", True, False),
]
for ri, (id_, name, price, avail, is_true, highlight) in enumerate(rows_data):
    ry = Inches(2.55) + ri*Inches(0.5)
    row_bg = DB_BG
    row_b = box(s, Inches(0.35), ry, Inches(9.3), Inches(0.48), fill=row_bg)
    row_b.line.color.rgb = DB_MID; row_b.line.width = Pt(0.5)
    tf_row = row_b.text_frame; tf_row.word_wrap = False
    p_row = tf_row.paragraphs[0]
    for cell_txt, cell_color in [(id_.ljust(8), DB_TEXT), (name.ljust(30), DB_TEXT), (price.ljust(12), DB_TEXT), (avail, DB_GREEN if is_true else DB_RED)]:
        r_row = p_row.add_run(); r_row.text = cell_txt
        r_row.font.size = Pt(10); r_row.font.color.rgb = cell_color
        r_row.font.bold = False

txt(s, "☝️  คอลัมน์ is_available คือตัวควบคุม —  true = แสดงในเมนู  ·  false = ซ่อนจากลูกค้า",
    Inches(0.35), Inches(4.65), Inches(9.3), Inches(0.5),
    size=Pt(11), bold=True, color=RED)
tip_box(s, "ℹ️ is_available แปลตรง ๆ ว่า 'มีพร้อมสั่งไหม?'",
        "true = มี (ลูกค้าเห็น)  ·  false = ไม่มี (ลูกค้าไม่เห็น แต่ข้อมูลยังอยู่)",
        Inches(0.35), Inches(5.3), Inches(9.3), Inches(0.9))

# ════════════════════════════════════════════════════════════
#  SLIDE 10 — คลิกเปลี่ยนค่า
# ════════════════════════════════════════════════════════════
s = slide(); slide_header(s, "ซ่อนเมนู · ขั้น 3", "คลิกเซลล์ → พิมพ์ false → Enter")
bg(s, CREAM)
txt(s, "ตัวอย่าง: ต้องการซ่อน 'ชุดไส้พะโล้' เพราะขายหมด",
    Inches(0.4), Inches(1.3), Inches(9.2), Inches(0.35), size=Pt(11), color=MID)

# DB mockup condensed
db_bar2 = box(s, Inches(0.35), Inches(1.7), Inches(9.3), Inches(0.4), fill=DB_BG)
db_bar2.line.fill.background()
tf_db2 = db_bar2.text_frame
p_db2 = tf_db2.paragraphs[0]
for txt_p, clr in [("⚡ Supabase", DB_GREEN), ("  ›  Table Editor › menu", DB_MUTED)]:
    rr = p_db2.add_run(); rr.text = txt_p
    rr.font.size = Pt(10); rr.font.bold = (clr == DB_GREEN); rr.font.color.rgb = clr

thead2 = box(s, Inches(0.35), Inches(2.1), Inches(9.3), Inches(0.35), fill=DB_MID)
thead2.line.fill.background()
tf_th2 = thead2.text_frame; p_th2 = tf_th2.paragraphs[0]
for col_name2 in ["id        ", "name (ชื่อเมนู)               ", "price     ", "is_available"]:
    r_th2 = p_th2.add_run(); r_th2.text = col_name2
    r_th2.font.size = Pt(9); r_th2.font.color.rgb = DB_MUTED; r_th2.font.bold = True

rows10 = [
    ("1 ","ขาหมูพะโล้เตาถ่าน", "180", "✓ true", DB_GREEN, False),
    ("5 ","ชุดไส้พะโล้เตาถ่าน", "120", "false|", DB_RED, True),
    ("4 ","ข้าวขาหมูพะโล้", "60", "✓ true", DB_GREEN, False),
]
for ri2, (id2, name2, price2, avail2, av_clr, hl) in enumerate(rows10):
    ry2 = Inches(2.45) + ri2*Inches(0.55)
    rbg = RGBColor(0x0D,0x21,0x16) if hl else DB_BG
    rb = box(s, Inches(0.35), ry2, Inches(9.3), Inches(0.52), fill=rbg)
    if hl:
        rb.line.color.rgb = DB_GREEN; rb.line.width = Pt(2)
    else:
        rb.line.color.rgb = DB_MID; rb.line.width = Pt(0.5)
    tf_r = rb.text_frame; p_r = tf_r.paragraphs[0]
    for ct, cc, cbold in [(id2.ljust(8), DB_TEXT, False),
                          (name2.ljust(30), CODE_YEL if hl else DB_TEXT, hl),
                          (price2.ljust(12), DB_TEXT, False),
                          (avail2, av_clr, True)]:
        rr2 = p_r.add_run(); rr2.text = ct
        rr2.font.size = Pt(10); rr2.font.color.rgb = cc; rr2.font.bold = cbold

txt(s, "⬆️  แถวนี้กำลังแก้ไข — พิมพ์  false  แล้วกด  Enter",
    Inches(0.35), Inches(4.08), Inches(9.3), Inches(0.4),
    size=Pt(11), bold=True, color=RED)

steps10 = [
    ("①","ดับเบิลคลิกที่เซลล์ is_available ของเมนูที่ต้องการซ่อน","เซลล์จะเปลี่ยนเป็นโหมดแก้ไข (เห็น cursor กะพริบ)"),
    ("②","ลบค่าเดิม พิมพ์ false (ตัวพิมพ์เล็กทั้งหมด) → Enter","เมนูหายจากแอปลูกค้าทันที ไม่ต้อง Refresh"),
]
for i10, (n10, t10, b10) in enumerate(steps10):
    step_row(s, n10, t10, b10, Inches(4.65) + i10*Inches(1.0))

# ════════════════════════════════════════════════════════════
#  SLIDE 11 — Before / After
# ════════════════════════════════════════════════════════════
s = slide(); slide_header(s, "ซ่อนเมนู · ผลลัพธ์", "ลูกค้าเห็นอะไรก่อน/หลังซ่อน")

# Before box
bb = box(s, Inches(0.35), Inches(1.4), Inches(4.55), Inches(3.8), fill=WHITE,
         border=RGBColor(0xFE,0xCA,0xCA))
box(s, Inches(0.35), Inches(1.4), Inches(4.55), Inches(0.45), fill=RGBColor(0xFE,0xE2,0xE2))
txt(s, "ก่อนซ่อน ❌", Inches(0.35), Inches(1.4), Inches(4.55), Inches(0.45),
    size=Pt(11), bold=True, color=RGBColor(0x99,0x1B,0x1B), align=PP_ALIGN.CENTER)
before_items = [("🍖 ขาหมูพะโล้","180฿"),("🫀 ชุดไส้พะโล้","120฿"),
                ("🍚 ข้าวขาหมู","60฿"),("☕ กาแฟสด","65฿")]
for bi, (nm, pr) in enumerate(before_items):
    by = Inches(1.95) + bi*Inches(0.62)
    bg_row = RGBColor(0xFE,0xF9,0xC3) if bi==1 else WHITE
    box(s, Inches(0.45), by, Inches(4.35), Inches(0.55), fill=bg_row)
    txt(s, nm, Inches(0.55), by+Inches(0.1), Inches(2.8), Inches(0.38),
        size=Pt(11), bold=(bi==1), color=DARK)
    txt(s, pr, Inches(3.8), by+Inches(0.1), Inches(0.9), Inches(0.38),
        size=Pt(11), bold=(bi==1), color=DARK, align=PP_ALIGN.RIGHT)

# After box
ab = box(s, Inches(5.1), Inches(1.4), Inches(4.55), Inches(3.8), fill=WHITE,
         border=RGBColor(0x86,0xEF,0xAC))
box(s, Inches(5.1), Inches(1.4), Inches(4.55), Inches(0.45), fill=RGBColor(0xDC,0xFC,0xE7))
txt(s, "หลังซ่อน ✅", Inches(5.1), Inches(1.4), Inches(4.55), Inches(0.45),
    size=Pt(11), bold=True, color=GREEN_GOOD, align=PP_ALIGN.CENTER)
after_items = [("🍖 ขาหมูพะโล้","180฿"),("🍚 ข้าวขาหมู","60฿"),("☕ กาแฟสด","65฿")]
for ai, (nm, pr) in enumerate(after_items):
    ay = Inches(1.95) + ai*Inches(0.62)
    box(s, Inches(5.2), ay, Inches(4.35), Inches(0.55), fill=WHITE)
    txt(s, nm, Inches(5.3), ay+Inches(0.1), Inches(2.8), Inches(0.38), size=Pt(11), color=DARK)
    txt(s, pr, Inches(8.5), ay+Inches(0.1), Inches(0.9), Inches(0.38),
        size=Pt(11), color=DARK, align=PP_ALIGN.RIGHT)
txt(s, "✓ 'ชุดไส้' หายออกไปแล้ว", Inches(5.2), Inches(3.83), Inches(4.35), Inches(0.35),
    size=Pt(10), bold=True, color=RGBColor(0x16,0xA3,0x4A))

tip_box(s, "↩️ อยากเปิดกลับมา?",
        "กลับไป Supabase → ดับเบิลคลิกเซลล์เดิม → พิมพ์ true → Enter  เมนูกลับมาทันที",
        Inches(0.35), Inches(5.35), Inches(9.3), Inches(0.85),
        bg_color=WARN_BG, title_color=WARN_FG, border_color=RGBColor(0xCC,0xA0,0x10))
tip_box(s, "✅ ซ่อนไม่ลบ — ข้อมูลไม่หายไปไหน",
        "รูปภาพ ราคา ตัวเลือก ทุกอย่างยังอยู่ครบ แค่ลูกค้าไม่เห็นชั่วคราว",
        Inches(0.35), Inches(6.3), Inches(9.3), Inches(0.85),
        bg_color=RGBColor(0xF0,0xFD,0xF4), title_color=GREEN_GOOD,
        border_color=RGBColor(0x22,0xC5,0x5E))

# ════════════════════════════════════════════════════════════
#  SLIDE 12 — SQL เร็วกว่า
# ════════════════════════════════════════════════════════════
s = slide(); slide_header(s, "💡 ข้อเสนอแนะ", "วิธีเร็วกว่า — ใช้ SQL แทน")
tip_box(s, "🎯 ทำไม SQL ถึงง่ายกว่า Table Editor?",
        "Table Editor: ต้องหาแถว → เลื่อนหา column → คลิก → แก้\nSQL: แค่พิมพ์ชื่อเมนู → กด Run → เสร็จ ใช้เวลาไม่ถึง 10 วินาที",
        Inches(0.35), Inches(1.4), Inches(9.3), Inches(1.1),
        bg_color=HINT_BG, title_color=HINT_FG, border_color=RGBColor(0x63,0x66,0xF1))
txt(s, "เมนูซ้าย → SQL Editor → New Query → วางโค้ดนี้ → กด Run",
    Inches(0.35), Inches(2.65), Inches(9.3), Inches(0.35), size=Pt(11), color=MID)

# SQL code box
sql_bg = box(s, Inches(0.35), Inches(3.05), Inches(9.3), Inches(2.1), fill=CODE_BG)
sql_bg.line.fill.background()
tf_sql = sql_bg.text_frame; tf_sql.word_wrap = False
sql_lines = [
    [("-- 🙈 ซ่อนเมนู (เปลี่ยนแค่ชื่อเมนู)", CODE_GRAY, False)],
    [("UPDATE ", CODE_PURP, True), ("public.menu", CODE_BLUE, False)],
    [("SET ", CODE_PURP, True), ("is_available = ", CODE_BLUE, False), ("false", CODE_GREEN, True)],
    [("WHERE ", CODE_PURP, True), ("name ", CODE_BLUE, False), ("ILIKE ", CODE_PURP, True), ("'%ชุดไส้%';", CODE_GREEN, False)],
    [("", CODE_GRAY, False)],
    [("-- 👁️ เปิดกลับมา", CODE_GRAY, False)],
    [("UPDATE ", CODE_PURP, True), ("public.menu", CODE_BLUE, False)],
    [("SET ", CODE_PURP, True), ("is_available = ", CODE_BLUE, False), ("true", CODE_GREEN, True)],
    [("WHERE ", CODE_PURP, True), ("name ", CODE_BLUE, False), ("ILIKE ", CODE_PURP, True), ("'%ชุดไส้%';", CODE_GREEN, False)],
]
first_sql = True
for line_parts in sql_lines:
    if first_sql:
        p_sql = tf_sql.paragraphs[0]; first_sql = False
    else:
        p_sql = tf_sql.add_paragraph()
    p_sql.alignment = PP_ALIGN.LEFT
    for part_txt, part_clr, part_bold in line_parts:
        r_sql = p_sql.add_run(); r_sql.text = part_txt
        r_sql.font.size = Pt(10); r_sql.font.color.rgb = part_clr
        r_sql.font.bold = part_bold; r_sql.font.name = "Courier New"

tip_box(s, "✏️ เปลี่ยนแค่ 'ชุดไส้' = ชื่อ (หรือส่วนของชื่อ) เมนูที่ต้องการซ่อน",
        "เช่น ถ้าต้องการซ่อน 'ขาหมูพะโล้เตาถ่าน' ก็พิมพ์  ขาหมู",
        Inches(0.35), Inches(5.3), Inches(9.3), Inches(0.85),
        bg_color=HINT_BG, title_color=HINT_FG, border_color=RGBColor(0x63,0x66,0xF1))
tip_box(s, "⚠️ ILIKE คืออะไร?",
        "คือค้นหาชื่อที่มีคำนั้นอยู่ เช่น '%ไส้%' จับทุกเมนูที่มีคำว่า 'ไส้' ในชื่อ",
        Inches(0.35), Inches(6.28), Inches(9.3), Inches(0.9),
        bg_color=WARN_BG, title_color=WARN_FG, border_color=RGBColor(0xCC,0xA0,0x10))

# ════════════════════════════════════════════════════════════
#  SLIDE 13 — แก้ราคา
# ════════════════════════════════════════════════════════════
s = slide(); slide_header(s, "แก้ราคา", "แก้ราคาเมนูด้วยตัวเอง")
steps13 = [
    ("1","Supabase → Table Editor → ตาราง menu","เห็นรายการเมนูทั้งหมด"),
    ("2","หาเมนูที่ต้องการ → ดับเบิลคลิกเซลล์ 'price'","คลิกที่ตัวเลขราคา เซลล์จะเข้าโหมดแก้ไข"),
    ("3","พิมพ์ราคาใหม่ → กด Enter","แอปลูกค้าอัปเดตภายใน 1–2 วินาที ไม่ต้อง deploy ใหม่"),
]
for i13, (n13, t13, b13) in enumerate(steps13):
    step_row(s, n13, t13, b13, Inches(1.4) + i13*Inches(1.0))

# mini DB mockup
db13 = box(s, Inches(0.35), Inches(4.65), Inches(9.3), Inches(0.4), fill=DB_BG)
db13.line.fill.background()
tf13 = db13.text_frame; p13 = tf13.paragraphs[0]
for tp, tc in [("⚡ Supabase", DB_GREEN), ("  ›  Table Editor › menu", DB_MUTED)]:
    r13 = p13.add_run(); r13.text = tp; r13.font.size=Pt(10); r13.font.color.rgb=tc; r13.font.bold=(tc==DB_GREEN)

th13 = box(s, Inches(0.35), Inches(5.05), Inches(9.3), Inches(0.35), fill=DB_MID)
th13.line.fill.background()
tf_th13 = th13.text_frame; p_th13 = tf_th13.paragraphs[0]
for cn in ["id        ","name                          ","price          ","is_available"]:
    r_th13 = p_th13.add_run(); r_th13.text = cn
    r_th13.font.size = Pt(9); r_th13.font.color.rgb = DB_MUTED; r_th13.font.bold = True

rows13 = [("1 ","ขาหมูพะโล้เตาถ่าน","200|",True),("4 ","ข้าวขาหมูพะโล้","60 ",False)]
for ri13, (id13, nm13, pr13, hl13) in enumerate(rows13):
    ry13 = Inches(5.4) + ri13*Inches(0.52)
    rb13 = box(s, Inches(0.35), ry13, Inches(9.3), Inches(0.5),
               fill=RGBColor(0x0D,0x21,0x16) if hl13 else DB_BG)
    rb13.line.color.rgb = DB_GREEN if hl13 else DB_MID
    rb13.line.width = Pt(2 if hl13 else 0.5)
    tf_r13 = rb13.text_frame; p_r13 = tf_r13.paragraphs[0]
    for ct13, cc13 in [(id13.ljust(8), DB_TEXT),(nm13.ljust(32), DB_TEXT),
                       (pr13.ljust(16), CODE_YEL if hl13 else DB_TEXT),("✓ true", DB_GREEN)]:
        rr13 = p_r13.add_run(); rr13.text = ct13
        rr13.font.size = Pt(10); rr13.font.color.rgb = cc13; rr13.font.bold = hl13 and ct13.strip() in ["200|"]

txt(s, "☝️  คลิกเซลล์ price → พิมพ์ราคาใหม่ (200) → กด Enter เสร็จ",
    Inches(0.35), Inches(6.55), Inches(9.3), Inches(0.4),
    size=Pt(11), bold=True, color=RED)

# ════════════════════════════════════════════════════════════
#  SLIDE 14 — Supabase Paused
# ════════════════════════════════════════════════════════════
s = slide(); slide_header(s, "แก้ปัญหา", "Supabase ขึ้น 'Paused' — ทำอย่างไร?")
tip_box(s, "⚠️ Paused เกิดจากอะไร?",
        "Supabase แบบฟรีจะหยุดทำงานอัตโนมัติถ้าไม่ได้ใช้งาน 7 วัน\nเมนูลูกค้ายังดูได้ แต่การแก้ราคาจะยังไม่มีผล",
        Inches(0.35), Inches(1.4), Inches(9.3), Inches(1.05),
        bg_color=WARN_BG, title_color=WARN_FG, border_color=RGBColor(0xCC,0xA0,0x10))
steps14 = [
    ("1","เข้า supabase.com → Sign in → คลิกโปรเจกต์","จะเห็นสถานะ 'Paused' หรือ 'Unhealthy'"),
    ("2","กดปุ่มสีเขียว 'Resume project'","โปรเจกต์จะเริ่มต้นใหม่ ใช้เวลา 1–2 นาที"),
    ("3","รอ 2 นาที → กด Refresh หน้าแอปเมนู","ทุกอย่างกลับมาปกติ ราคาอัปเดตได้ตามเดิม"),
]
for i14, (n14, t14, b14) in enumerate(steps14):
    step_row(s, n14, t14, b14, Inches(2.65) + i14*Inches(1.0))
tip_box(s, "💡 ป้องกัน Paused — เข้าสัปดาห์ละครั้ง",
        "แค่เปิดหน้า Supabase Table Editor ก็นับว่า active แล้ว ไม่ต้องทำอะไรเพิ่ม",
        Inches(0.35), Inches(5.85), Inches(9.3), Inches(0.85),
        bg_color=HINT_BG, title_color=HINT_FG, border_color=RGBColor(0x63,0x66,0xF1))
tip_box(s, "📌 ถ้าร้านเปิดทุกวัน แนะนำ Supabase Pro",
        "$25/เดือน — ไม่มี pause อัตโนมัติ ไม่ต้องกังวลเรื่องนี้อีก",
        Inches(0.35), Inches(6.8), Inches(9.3), Inches(0.55))

# ════════════════════════════════════════════════════════════
#  SLIDE 15 — จบ + Zoom
# ════════════════════════════════════════════════════════════
s = slide(); bg(s, SAGE)
txt(s, "🎥", Inches(4.3), Inches(0.9), Inches(1.4), Inches(0.9),
    size=Pt(48), color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "จบสไลด์แล้ว!", Inches(1), Inches(1.85), Inches(8), Inches(0.8),
    size=Pt(36), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "ดูครบแล้ว ถ้ามีคำถามอะไรเพิ่ม ทักเราทาง LINE ได้เลยครับ",
    Inches(1), Inches(2.7), Inches(8), Inches(0.5),
    size=Pt(14), color=WHITE, align=PP_ALIGN.CENTER)

# Zoom CTA box
zoom = box(s, Inches(2.0), Inches(3.4), Inches(6), Inches(3.6),
           fill=RGBColor(0xFF,0xFF,0xFF))
zoom.line.color.rgb = RGBColor(0xFF,0xFF,0xFF)
zoom.fill.fore_color.rgb = RGBColor(0xFF,0xFF,0xFF)
tf_z = zoom.text_frame; tf_z.word_wrap = True
p_z0 = tf_z.paragraphs[0]; p_z0.alignment = PP_ALIGN.LEFT
r_z0 = p_z0.add_run(); r_z0.text = "📅 อย่าลืมนัด Zoom สอนงานจริง 30 นาที"
r_z0.font.size = Pt(13); r_z0.font.bold = True; r_z0.font.color.rgb = DARK
zoom_items = [
    "✓  ทดลองซ่อนเมนูพร้อมกัน live",
    "✓  ทดลองแก้ราคาใน Supabase",
    "✓  ดูออเดอร์เข้า LINE OA จริง",
    "✓  ตอบคำถามที่ค้างอยู่ทุกข้อ",
    "✓  บันทึก Zoom ไว้ดูย้อนหลังได้",
]
for zi in zoom_items:
    p_zi = tf_z.add_paragraph(); p_zi.alignment = PP_ALIGN.LEFT
    r_zi = p_zi.add_run(); r_zi.text = zi
    r_zi.font.size = Pt(12); r_zi.font.color.rgb = DARK

# ── Save ────────────────────────────────────────────────────
out = "/home/user/khaomoo-menu/docs/slides-premium.pptx"
prs.save(out)
print("Saved:", out)
